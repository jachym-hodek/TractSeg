#!/usr/bin/env python3
"""
png_to_presentation.py

Usage:
  python3 png_to_presentation.py DIR1 [DIR2 ...] [OUTPUT.pptx]

Rules:
- Sections = immediate subdirectory names. If DIR1/foo and DIR2/foo exist,
  their images are merged into one "foo" section.
- Within a section, images are matched by *filename*. e.g., foo/img1.png from
  different DIRs go on the same slide.
- Images that exist in only one DIR still get their own slide.
- Root-level images (directly under DIR1, DIR2, ...) are NOT merged across DIRs.
  They appear in sections named "<DIR basename> (root)".

Notes:
- Adds a divider slide per section.
- Natural-sort order for filenames (e.g., 1, 2, 10).
- Handles 1..N images per slide with an auto grid. Adds per-image captions
  indicating which top-level DIR it came from.
"""

import sys, os, glob, re, math
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

EMU_PER_INCH = 914400

def natural_key(s: str):
    b = os.path.basename(s)
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r'(\d+)', b)]

def parse_args(argv):
    if len(argv) < 2:
        sys.exit("Usage: python3 png_to_presentation.py DIR1 [DIR2 ...] [OUTPUT.pptx]")
    args = argv[1:]
    out = "merged_images.pptx"
    if len(args) >= 2 and args[-1].lower().endswith(".pptx"):
        out = args[-1]
        roots = args[:-1]
    else:
        roots = args
    if not roots:
        sys.exit("Need at least one input directory.")
    bad = [p for p in roots if not os.path.isdir(p)]
    if bad:
        sys.exit("Not directories: " + ", ".join(bad))
    return roots, out

def unique_root_labels(roots):
    """Return list of (root_path, display_label) where duplicate basenames get suffixes."""
    bases = [os.path.basename(os.path.normpath(r)) or r for r in roots]
    counts = Counter(bases)
    seen = Counter()
    labels = []
    for r, b in zip(roots, bases):
        seen[b] += 1
        if counts[b] > 1:
            labels.append((r, f"{b} ({seen[b]})"))
        else:
            labels.append((r, b))
    return labels

def gather_sections(roots, root_labels):
    """
    Build:
      sections[name][filename] -> list of entries for different roots that have that file
      where entry = (root_index, root_label, full_path)
    Also create per-root 'root' sections named "<root_label> (root)" for images in the root.
    """
    sections = defaultdict(lambda: defaultdict(list))

    # Map root path -> (index, label)
    index_of = {rp: i for i, (rp, _) in enumerate(root_labels)}
    label_of = {rp: lb for rp, lb in root_labels}

    for rp, lb in root_labels:
        ridx = index_of[rp]

        # Root-level PNGs -> a private section that does NOT merge across roots
        root_imgs = sorted([*glob.glob(os.path.join(rp, "*.png")),
                            *glob.glob(os.path.join(rp, "*.PNG"))], key=natural_key)
        if root_imgs:
            root_section = f"{lb} (root)"
            for p in root_imgs:
                fname = os.path.basename(p)
                sections[root_section][fname].append((ridx, lb, p))

        # Immediate subdirectories only
        for name in sorted(os.listdir(rp), key=str.lower):
            sub = os.path.join(rp, name)
            if not os.path.isdir(sub):
                continue
            imgs = sorted([*glob.glob(os.path.join(sub, "*.png")),
                           *glob.glob(os.path.join(sub, "*.PNG"))], key=natural_key)
            if not imgs:
                continue
            sec = os.path.basename(sub)  # merged across roots by this name
            for p in imgs:
                fname = os.path.basename(p)
                sections[sec][fname].append((ridx, lb, p))

    # Within each filename group, order entries by the order the roots were provided
    for sec, files in sections.items():
        for fname, entries in files.items():
            files[fname] = sorted(entries, key=lambda e: e[0])

    return sections

def add_section_divider(prs, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left = top = Inches(1)
    width = prs.slide_width - Inches(2)
    height = prs.slide_height - Inches(2)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    p.alignment = PP_ALIGN.CENTER
    run.font.bold = True
    run.font.size = Pt(60)

def add_title(slide, text: str, prs):
    top = Inches(0.3)
    margin = Inches(0.5)
    width = prs.slide_width - 2 * margin
    box = slide.shapes.add_textbox(margin, top, width, Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    p.alignment = PP_ALIGN.CENTER
    r.font.size = Pt(24)
    r.font.bold = True

def add_multi_image_slide(prs, title_text: str, items):
    """
    items: list of (root_label, image_path) already ordered by root index.
    Layout: auto grid with captions under each cell.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title_text, prs)

    n = len(items)
    if n == 0:
        return

    # Page geometry
    margin = Inches(0.5)
    top_after_title = Inches(1.1)
    gap = Inches(0.25)
    caption_h = Inches(0.3)

    avail_w = prs.slide_width - 2 * margin
    avail_h = prs.slide_height - top_after_title - margin

    # Grid: try to be square-ish
    cols = max(1, math.ceil(math.sqrt(n)))
    rows = math.ceil(n / cols)

    # If a single row would be too short (wide images), prefer more columns up to n
    while (rows - 1) * cols >= n:
        rows -= 1

    cell_w = (avail_w - (cols - 1) * gap) / cols
    cell_h = (avail_h - (rows - 1) * gap) / rows

    # Place each image into its cell
    for idx, (label, path) in enumerate(items):
        r = idx // cols
        c = idx % cols
        left = margin + c * (cell_w + gap)
        top = top_after_title + r * (cell_h + gap)

        # Fit image preserving aspect ratio inside (cell_w, cell_h - caption_h)
        with Image.open(path) as im:
            iw, ih = im.width, im.height
        # Avoid div-by-zero
        iw = iw or 1
        ih = ih or 1
        target_w = cell_w
        target_h = cell_h - caption_h
        if target_h <= 0:
            target_h = cell_h

        ratio = min(float(target_w) / (iw), float(target_h) / (ih))
        pic_w = int(iw * ratio)
        pic_h = int(ih * ratio)

        # Center inside cell
        pic_left = int(left + (cell_w - pic_w) / 2)
        pic_top = int(top + (cell_h - caption_h - pic_h) / 2)

        slide.shapes.add_picture(path, pic_left, pic_top, width=pic_w, height=pic_h)

        # Caption under image
        cap_left = left
        cap_top = top + cell_h - caption_h
        cap_box = slide.shapes.add_textbox(cap_left, cap_top, cell_w, caption_h)
        tf = cap_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        p.alignment = PP_ALIGN.CENTER
        run.font.size = Pt(12)

def build_deck(sections, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Order sections: alphabetically, but keep any "(root)" sections last so real names come first
    named = sorted([s for s in sections.keys() if not s.endswith(" (root)")], key=str.lower)
    roots = sorted([s for s in sections.keys() if s.endswith(" (root)")], key=str.lower)
    ordered_sections = named + roots

    total_image_groups = 0
    for sec in ordered_sections:
        add_section_divider(prs, sec)
        files = sections[sec]
        for fname in sorted(files.keys(), key=natural_key):
            entries = files[fname]  # list of (root_index, root_label, path)
            # Already ordered by root index; collapse to (label, path)
            items = [(label, path) for (_idx, label, path) in entries]
            title = f"{sec} — {fname}"
            add_multi_image_slide(prs, title, items)
            total_image_groups += 1

    prs.save(out_path)
    return len(ordered_sections), total_image_groups

def main():
    roots, out = parse_args(sys.argv)
    # Normalize roots and produce stable display labels
    roots = [os.path.abspath(r) for r in roots]
    root_labels = unique_root_labels(roots)  # [(root_path, display_label)]
    sections = gather_sections(roots, root_labels)

    if not sections:
        sys.exit("No PNGs found under the provided directories.")

    n_sections, n_groups = build_deck(sections, out)
    print(f"Wrote {out}: {n_groups} slides across {n_sections} sections.")

if __name__ == "__main__":
    main()
